from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from os.path import exists
import re
from PIL import Image

class Deck:
   BAD_VALUE = -999

   def __init__(self, fnpres):
      self.fnpres = fnpres
      self.slides = []
      self.slide_los = []

      if (not isinstance(self.fnpres, str)):
         raise TypeError("--- presentation must be a PowerPoint filename string ---")
      else:
         if (not exists(self.fnpres)):
            raise FileNotFoundError("--- presentation file does not exist ---")

      #  Instantiate the presentation
      self.pres = Presentation(self.fnpres)

      #  Read the slide layout components into a list of dictionaries
      i = 0
      while (i < len(self.pres.slide_layouts)):
         self.phDict = {}
         for shape in self.pres.slide_layouts[i].placeholders:
            self.phDict[shape.name] = shape.placeholder_format.idx
         self.slide_los.append(self.phDict)
         i += 1

   #  Print deck description
   def __str__(self):
      return f"Presentation file: {self.fnpres}"

   #  Save the deck
   def save(self, fn):
      if (not isinstance(fn, str)):
         raise TypeError("--- filename must be a string to save presentation ---")
      for s in range(0, len(self.slides)):
         self.render_slide(self.slides[s])
      self.pres.save(fn)
      return True

   #  Show deck filename
   def show_filename(self):
      print("Deck file: {}".format(self.fnpres))
      return True

   #  Print number of slides in the deck
   def numslides(self):
      return len(self.slides)

   #  Add a Slide object to the deck
   def add_slide(self, slide, index=None):
      #  Check slide is a Slide object
      if (not isinstance(slide, Slide)):
         raise TypeError("--- can only add Slide objects ---")

      #  Determine appropriate layout using find_layout() method
      lolist = self.find_layout(slide)
      lo = lolist[0]

      if (lo == self.BAD_VALUE):
         raise ValueError("--- couldn't find conforming layout ---")

      #  If specified, ensure index is within the range of slides and then
      #  insert; otherwise, append it
      if (index is not None):
         if (index < len(self.slides)):
            self.slides.insert(index, slide)
      else:
         self.slides.append(slide)

      return True

   #  Render Slide object(s) to the deck
   def render_slide(self, slide, layout=None, index=None):
      #  Check slide is a Slide object
      if (not isinstance(slide, Slide)):
         raise TypeError("--- can only add Slide objects ---")

      #  If specified, ensure layout is an integer within the range of layouts,
      #  else determine appropriate layout using find_layout() method
      if (layout is not None):
         if (isinstance(layout, int)):
            if (layout >= len(self.pres.slide_layouts)):
               raise IndexError("--- layout outside range of slide layouts ---")
            lo = layout
         else:
            raise TypeError("--- layout reference must be an integer ---")
      else:
         lolist = self.find_layout(slide)
         lo = lolist[0]

      if (lo == self.BAD_VALUE):
         print("---  for {}:".format(slide.name))
         raise ValueError("--- couldn't find conforming layout ---")

      #  If specified, ensure index is within the range of slides and then
      #  insert; otherwise, append it
      if (index is not None):
         if (index < len(self.slides)):
            self.slides.insert(index, slide)
      else:
         self.slides.append(slide)
      new_slide = self.pres.slides.add_slide(self.pres.slide_layouts[lo])

      #  Fill the title placeholder, if applicable
      if (slide.title is not None):
         ph = self.get_title_ph(lo)
         new_slide.placeholders[ph.pop()].text = slide.title

      #  Fill the slide exhibit (Picture) placeholders, if applicable
      if (slide.num_exhibits() > 0):
         ph = self.get_picture_ph(lo)
         i = 0
         while (i < len(slide.exhibits)):
            new_slide.placeholders[ph.pop()].insert_picture(slide.exhibits[i])
            i += 1

      #  Fill the Main text box with bullets, if applicable
      print("Rendering............")
      if (slide.num_main_bullets() > 0):
         #  Set up the text box
         ph = self.get_main_ph(lo)
         tf_main = new_slide.shapes.placeholders[ph.pop()].text_frame
         tf_main.clear()
         tf_main.margin_left = 0
         tf_main.vertical_anchor = MSO_ANCHOR.TOP
         tf_main.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
         #  Loop through and implement the formatting "runs"
         for i in range(0, len(slide.run_main)):
            if (re.search("ParaLevel", slide.run_main[i])):
               (junk, level) = slide.run_main[i].split("Level")
               p = tf_main.add_paragraph()
               p.level = int(level)
               p.alignment = PP_ALIGN.LEFT
            elif (re.search("::", slide.run_main[i])):
               (fnm, fsz, bbold, bitalic) = slide.run_main[i].split("::")
               run = p.add_run()
               run.font.name = fnm
               run.font.size = Pt(int(fsz))
               if (bbold == "True"):
                  run.font.bold = True
               if (bitalic == "True"):
                  run.font.italic = True
            else:
               run.text = slide.run_main[i]

      #  Fill the Margin text box with bullets, if applicable
      if (slide.num_margin_bullets() > 0):
         #  Set up the text box
         ph = self.get_margin_ph(lo)
         tf_margin = new_slide.shapes.placeholders[ph.pop()].text_frame
         tf_margin.clear()
         tf_margin.margin_left = 0
         tf_margin.vertical_anchor = MSO_ANCHOR.TOP
         tf_margin.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
         #  Loop through and implement the formatting "runs"
         for i in range(0, len(slide.run_marg)):
            if (re.search("ParaLevel", slide.run_marg[i])):
               (junk, level) = slide.run_marg[i].split("Level")
               p = tf_margin.add_paragraph()
               p.level = int(level)
               p.alignment = PP_ALIGN.LEFT
            elif (re.search("::", slide.run_marg[i])):
               (fnm, fsz, bbold, bitalic) = slide.run_marg[i].split("::")
               run = p.add_run()
               run.font.name = fnm
               run.font.size = Pt(int(fsz))
               if (bbold == "True"):
                  run.font.bold = True
               if (bitalic == "True"):
                  run.font.italic = True
            else:
               run.text = slide.run_marg[i]

      #  Fill the footer with footnotes, if applicable
      if (slide.num_footnotes() > 0):
         #  Set up the text box
         ph = self.get_footer_ph(lo)
         tf_footer = new_slide.shapes.placeholders[ph.pop()].text_frame
         tf_footer.clear()
         tf_footer.margin_left = 0
         tf_footer.vertical_anchor = MSO_ANCHOR.TOP
         tf_footer.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
         #  Loop through and implement the formatting "runs"
         for i in range(0, len(slide.run_fn)):
            if (re.search("ParaLevel", slide.run_fn[i])):
               (junk, level) = slide.run_fn[i].split("Level")
               p = tf_footer.add_paragraph()
               p.level = int(level)
               p.alignment = PP_ALIGN.LEFT
            elif (re.search("::", slide.run_fn[i])):
               (fnm, fsz, bbold, bitalic) = slide.run_fn[i].split("::")
               run = p.add_run()
               run.font.name = fnm
               run.font.size = Pt(int(fsz))
               if (bbold == "True"):
                  run.font.bold = True
               if (bitalic == "True"):
                  run.font.italic = True
            else:
               run.text = slide.run_fn[i]

      return (len(self.slides) - 1)

   #  Delete a Slide object by name from the deck
   def del_slide(self, sn):
      numslb = len(self.slides)
      for i in range(0, len(self.slides)):
         if (self.slides[i].name == sn):
            self.slides.pop(i)
            break
      if (numslb == len(self.slides)):
         print("--- couldn't locate/delete slide named {} ---".format(sn))
         return False
      else:
         return True

   #  Show slides
   def show_slides(self):
      for s in self.slides:
         print(s.name)
      return True

   #  Find a suitable layout
   def find_layout(self, slide):
      #  Read available layouts from Deck object and select best fit based on
      #  active slide components
      ok_los = []
      i = 0
      while (i < len(self.slide_los)):
         c = True
         if (slide.title is not None):
            if (self.num_title_ph(i) == 0):
               c = False
         if (slide.num_exhibits() > 0):
            if (self.num_picture_ph(i) < slide.num_exhibits()):
               c = False
         if (slide.num_main_bullets() > 0):
            if (self.num_main_ph(i) == 0):
               c = False
         if (slide.num_margin_bullets() > 0):
            if (self.num_margin_ph(i) == 0):
               c = False
         if (slide.num_footnotes() > 0):
            if (self.num_footer_ph(i) == 0):
               c = False

         if (c == True):
            ok_los.append(i)
         i += 1

      if (len(ok_los) == 0):
         ok_los.append(self.BAD_VALUE)

      return ok_los

   #  Summarize slide layout placeholders
   def show_layouts(self):
      for lo in range(len(self.slide_los)):
         print('Layout {} placeholders:'.format(lo))
         print("----------------------------------------")
         for key in self.slide_los[lo].keys():
            print('{} ({})'.format(key, self.slide_los[lo][key]))
         print("----------------------------------------")
      return len(self.slide_los)

   #  Display placeholders in specified slide layout
   def show_placeholders(self, i):
      for key in self.slide_los[i].keys():
         print('{} ({})'.format(key, self.slide_los[i][key]))
      return len(self.slide_los[i])

   #  Return Picture placeholders in specified slide layout
   def get_picture_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Pic", key)):
            ph.append(self.slide_los[i][key])
      return ph

   #  Return Main placeholders in specified slide layout
   def get_main_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Main", key)):
            ph.append(self.slide_los[i][key])
      return ph

   #  Return Margin placeholders in specified slide layout
   def get_margin_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Margin", key)):
            ph.append(self.slide_los[i][key])
      return ph

   #  Return Footer placeholders in specified slide layout
   def get_footer_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Footer", key)):
            ph.append(self.slide_los[i][key])
      return ph

   #  Return Title placeholders in specified slide layout
   def get_title_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Title", key)):
            ph.append(self.slide_los[i][key])
      return ph

   #  Return Subtitle placeholders in specified slide layout
   def get_subtitle_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Subtitle", key)):
            ph.append(self.slide_los[i][key])
      return ph

   #  Return Date placeholders in specified slide layout
   def get_date_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Date", key)):
            ph.append(self.slide_los[i][key])
      return ph
 
   #  Return SlideNum placeholders in specified slide layout
   def get_slidenum_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^SlideNum", key)):
            ph.append(self.slide_los[i][key])
      return ph

  #  Return Table placeholders in specified slide layout
   def get_table_ph(self, i):
      ph = []
      for key in self.slide_los[i].keys():
         if (re.search("^Table", key)):
            ph.append(self.slide_los[i][key])
      return ph

   #  Return number of Picture placeholders in specified slide layout
   def num_picture_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Pic", key)):
            p += 1
      return p

   #  Return number of Main placeholders in specified slide layout
   def num_main_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Main", key)):
            p += 1
      return p

   #  Return number of Margin placeholders in specified slide layout
   def num_margin_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Margin", key)):
            p += 1
      return p

   #  Return number of Footer placeholders in specified slide layout
   def num_footer_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Footer", key)):
            p += 1
      return p

   #  Return number of Title placeholders in specified slide layout
   def num_title_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Title", key)):
            p += 1
      return p

   #  Return number of Subtitle placeholders in specified slide layout
   def num_subtitle_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Subtitle", key)):
            p += 1
      return p

   #  Return number of Date placeholders in specified slide layout
   def num_date_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Date", key)):
            p += 1
      return p

   #  Return number of SlideNum placeholders in specified slide layout
   def num_slidenum_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^SlideNum", key)):
            p += 1
      return p

   #  Return number of Table placeholders in specified slide layout
   def num_table_ph(self, i):
      p = 0
      for key in self.slide_los[i].keys():
         if (re.search("^Table", key)):
            p += 1
      return p

class Slide:
   def __init__(self, name, title=None, exhibits=None, bullets_main=None, bullets_marg=None, footnotes=None):
      self.name = name
      self.title = title
      self.exhibits = exhibits
      self.bullets_main = bullets_main
      self.bullets_marg = bullets_marg
      self.footnotes = footnotes

      #  Initialize arrays to hold Markdown-style "runs" for main and margin
      #  bullets and footnotes
      self.run_marg = []
      self.run_main = []
      self.run_fn = []

      if (self.title is not None):
         if (not isinstance(self.title, str)):
            raise TypeError("--- title must be a string ---")

      if (self.exhibits is not None):
         if (not isinstance(self.exhibits, list)):
            raise TypeError("--- exhibits must be a list ---")
         else:
            for e in self.exhibits:
               img = Image.open(e)
               fmt = img.get_format_mimetype()
               if ((fmt != "image/jpeg") & (fmt !="image/png")):
                  raise TypeError("--- exhibits must contain only PNG or JPEG images ---")

      if (self.bullets_main is not None):
         if (not isinstance(self.bullets_main, list)):
            raise TypeError("--- bullets_main must be a list ---")
         else:
            for b in self.bullets_main:
               if (not isinstance(b, str)):
                  raise TypeError("--- bullets_main must contain only strings ---")

      if (self.bullets_marg is not None):
         if (not isinstance(self.bullets_marg, list)):
            raise TypeError("--- bullets_marg must be a list ---")
         else:
            for b in self.bullets_marg:
               if (not isinstance(b, str)):
                  raise TypeError("--- bullets_marg must contain only strings ---")

      if (self.footnotes is not None):
         if (not isinstance(self.footnotes, list)):
            raise TypeError("--- footnotes must be a list ---")
         else:
            for f in self.footnotes:
               if (not isinstance(f, str)):
                  raise TypeError("--- footnotes must contain only strings ---")

   #  Print slide description
   def __str__(self):
      return f"{self.title}: exhibits={len(self.exhibits)}, main bullets={len(self.bullets_main)}, margin bullets={len(self.bullets_marg)}, footnotes={len(self.footnotes)}"

   #  Return number of exhibits
   def num_exhibits(self):
      if (self.exhibits is not None):
         return len(self.exhibits)
      else:
         return 0

   #  Return number of main bullets
   def num_main_bullets(self):
      if (self.bullets_main is not None):
         return len(self.bullets_main)
      else:
         return 0

   #  Return number of margin bullets
   def num_margin_bullets(self):
      if (self.bullets_marg is not None):
         return len(self.bullets_marg)
      else:
         return 0

   #  Return number of footnotes
   def num_footnotes(self):
      if (self.footnotes is not None):
         return len(self.footnotes)
      else:
         return 0

   #  Add title
   def add_title(self, title):
      if (isinstance(title, str)):
         self.title = title
         return True
      else:
         raise TypeError("--- title must be a string ---")
         return False

   #  Get title text
   def get_title(self):
      return self.title

   #  Add exhibit
   def add_exhibit(self, exhibit):
      #  Add exhibit to the exhibits list if it is a PNG or JPEG image
      img = Image.open(exhibit)
      fmt = img.get_format_mimetype()
      if ((fmt != "image/jpeg") & (fmt !="image/png")):
         raise TypeError("--- exhibits must be PNG or JPEG images ---")
      else:
         if (self.exhibits is None):
            self.exhibits = []
         self.exhibits.append(exhibit)
         return True

   #  Get list of exhibits
   def get_exhibits(self):
      return self.exhibits

   #  Add main bullet(s)
   def add_main_bullets(self, bullet):
      #  Add bullet(s) to the bullets_main list if it is a string object
      if (isinstance(bullet, str)):
         if (self.bullets_main is None):
            self.bullets_main = []
         self.parse_md(bullet)
         return True
      else:
         return False

   #  Get list of main bullets
   def get_main_bullets(self):
      return self.bullets_main

   #  Add margin bullet(s)
   def add_margin_bullets(self, bullet):
      #  Add bullet(s) to the bullets_margin list if it is a string object
      if (isinstance(bullet, str)):
         if (self.bullets_marg is None):
            self.bullets_marg = []
         self.parse_md(bullet)
         return True
      else:
         return False

   #  Get list of margin bullets
   def get_margin_bullets(self):
      return self.bullets_marg

   #  Add footnote(s)
   def add_footnotes(self, footnote):
      #  Add footnote to the footnotes list if it is a string object
      if (isinstance(footnote, str)):
         if (self.footnotes is None):
            self.footnotes = []
         self.parse_md(footnote)
         return True
      else:
         return False

   #  Get list of footnotes
   def get_footnotes(self):
      return self.footnotes

   #  Show exhibits
   def show_exhibits(self):
      if (self.exhibits is not None):
         for e in range(0, len(self.exhibits)):
            print(self.exhibits[e])
      else:
         print("--- No exhibits defined ---")
      return True

   #  Show main bullets
   def show_main_bullets(self):
      if (self.bullets_main is not None):
         for b in range(0, len(self.bullets_main)):
            print(self.bullets_main[b])
      else:
         print("--- No main bullets defined ---")
      return True

   #  Show margin bullets
   def show_margin_bullets(self):
      if (self.bullets_marg is not None):
         for b in range(0, len(self.bullets_marg)):
            print(self.bullets_marg[b])
      else:
         print("--- No margin bullets defined ---")
      return True

   #  Show footnotes
   def show_footnotes(self):
      if (self.footnotes is not None):
         for f in range(0, len(self.footnotes)):
            print(self.footnotes[f])
      else:
         print("--- No footnotes defined ---")
      return True

   #  Show main bullet runs
   def show_main_runs(self):
      if (self.run_main is not None and (len(self.run_main) > 0)):
         print("Main bullet runs:")
         print("----------------------------------------")
         for b in range(0, len(self.run_main)):
            print(self.run_main[b])
         print("----------------------------------------")
      else:
         print("--- No main bullet runs defined ---")
      return True

   #  Show margin bullet runs
   def show_margin_runs(self):
      if (self.run_marg is not None and (len(self.run_marg) > 0)):
         print("Margin bullet runs:")
         print("----------------------------------------")
         for b in range(0, len(self.run_marg)):
            print(self.run_marg[b])
         print("----------------------------------------")
      else:
         print("--- No margin bullet runs defined ---")
      return True

   #  Change title
   def chg_title(self, title):
      #  Change the slide title
      if (isinstance(title, str)):
         self.title = title
         return True
      else:
         return False

   #  Replace exhibit
   def rep_exhibit(self, index, exhibit):
      #  Check that the index is an integer within range of the exhibits list
      if (isinstance(index, int)):
         if (index >= len(exhibits)):
            raise IndexError("--- index out of range of exhibits list ---")
         else:
            return True
      else:
         raise TypeError("--- index must be an integer ---")

      #  Replace an exhibit in the list if it's a PNG or JPEG file
      img = Image.open(exhibit)
      fmt = img.get_format_mimetype()
      if ((fmt != "image/jpeg") & (fmt !="image/png")):
         raise TypeError("--- exhibits must be PNG or JPEG images ---")
      else:
         self.exhibits[index] = exhibit
         return True

   #  Replace main bullet
   def rep_main_bullet(self, index, mb):
      #  Check that the index is an integer within range of the main bullets list
      if (isinstance(index, int)):
         if (index >= len(self.bullets_main)):
            raise IndexError("--- index out of range of main bullets list ---")
         else:
            return True
      else:
         raise TypeError("--- index must be an integer ---")

      #  Replace a bullet in the list if it's a string
      if (isinstance(mb, str)):
         raise TypeError("--- main bullet must be a string ---")
      else:
         self.bullets_main[index] = mb
         return True

   #  Replace margin bullet
   def rep_margin_bullet(self, index, mb):
      #  Check that the index is an integer within range of the margin bullets list
      if (isinstance(index, int)):
         if (index >= len(self.bullets_margin)):
            raise IndexError("--- index out of range of margin bullets list ---")
         else:
            return True
      else:
         raise TypeError("--- index must be an integer ---")

      #  Replace a bullet in the list if it's a string
      if (isinstance(mb, str)):
         raise TypeError("--- margin bullet must be a string ---")
      else:
         self.bullets_margin[index] = mb
         return True

   #  Replace footnote
   def rep_footnote(self, index, fn):
      #  Check that the index is an integer within range of the footnotes list
      if (isinstance(index, int)):
         if (index >= len(self.footnotes)):
            raise IndexError("--- index out of range of footnotes list ---")
         else:
            return True
      else:
         raise TypeError("--- index must be an integer ---")

      #  Replace a footnote in the list if it's a string
      if (isinstance(fn, str)):
         raise TypeError("--- footnote must be a string ---")
      else:
         self.footnotes[index] = fn
         return True

   #  Parse a string for Markdown-style directives
   #  + (one or more) is a main bullet
   #  - (one or more) is a margin bullet
   #  ^ (one or more) is a footnote
   #  * italicizes text
   #  ** bolds text
   #  *# changes the font size
   #  *fontname changes the font
   def parse_md(self, s):
      #  Defaults
      dfnm = "Arial"
      dfsz = "11"
      dbold = "False"
      ditalic = "False"

      text = ""
      runtext = ""
      toks = re.split(' ', s)

      #  Check that a target symbol (+ or - or ^) is first, or throw an error
      if (re.search("^[\s]*[-+\^]+", toks[0])):
         if (re.search("^[\s]*[-]+", toks[0])):
            target = "margin"
         elif (re.search("^[\s]*[+]+", toks[0])):
            target = "main"
         else:
            target = "footnote"
      else:
         raise SyntaxError('--- must start with one or more "+" (main), "-" (margin) or "^" (footnote) ---')

      for j in range(0, len(toks)):
         #  Token is a margin or main bullet or footnote, so set a paragraph and
         #  determine the level
         if (re.search("^[\s]*[-+\^]+", toks[0])):
            #  Push any text to the storage arrays, then clear it
            if (text != ""):
               if (target == "margin"):
                  self.bullets_marg.append(text)
                  self.run_marg.append(runtext)
               elif (target == "main"):
                  self.bullets_main.append(text)
                  self.run_main.append(runtext)
               else:
                  self.footnotes.append(text)
                  self.run_fn.append(runtext)
               text = ""
               runtext = ""

            if (re.search("^[\s]*[-]+", toks[0])):
               target = "margin"
            elif (re.search("^[\s]*[+]+", toks[0])):
               target = "main"
            else:
               target = "footnote"

            #  (re)set the font name, size, bold and italics to defaults
            fnm = dfnm
            fsz = dfsz
            bold = dbold
            italic = ditalic

            #  Start a new paragraph and set the level (number of "-", "+" or "^"
            #  characters) and push the formatting strings to the target run array
            plev = str(len(toks[0]))
            if (target == "margin"):
               self.run_marg.append("--- ParaLevel" + plev)
               self.run_marg.append(fnm + "::" + fsz + "::" + bold + "::" + italic)
            elif (target == "main"):
               self.run_main.append("--- ParaLevel" + plev)
               self.run_main.append(fnm + "::" + fsz + "::" + bold + "::" + italic)
            else:
               self.run_fn.append("--- ParaLevel" + plev)
               self.run_fn.append(fnm + "::" + fsz + "::" + bold + "::" + italic)

         #  Token is a run directive, so interpret it
         elif (re.search("^[\*]+", toks[0])):
            #  If there's any bullet text, push it
            if (runtext != ""):
               if (target == "margin"):
                  self.run_marg.append(runtext)
               elif (target == "main"):
                  self.run_main.append(runtext)
               else:
                  self.run_fn.append(runtext)
               runtext = ""

            c = 0
            for i in range(0, len(toks[0])):
               if (toks[0][i] == "*"):
                  c += 1
            #  One asterisk indicates italics or font name or size
            if (c == 1):
               #  There's nothing after the asterisk, so toggle the italic setting
               if (len(toks[0]) == 1):
                  if (italic == "True"):
                     italic = "False"
                  else:
                     italic = "True"
               #  If there's a number after the asterisk, it's a font size
               elif (re.search("[0-9]+", toks[0][len(toks[0])-1])):
                  print("Parsing {}".format(toks[0]))
                  (ast, fsz) = toks[0].split('*')
               #  If there's a word after the asterisk, it's a font name
               elif (re.search("[A-Za-z ]+", toks[0][len(toks[0])-1])):
                  print("Parsing... {}".format(toks[0]))
                  (ast, fnm) = toks[0].split('*')
            #  Two asterisks indicate bold font, so toggle it
            elif (c == 2):
               if (bold == "True"):
                  bold = "False"
               else:
                  bold = "True"
            #  Set and push the text "decorations"
            if (target == "margin"):
               self.run_marg.append(fnm + "::" + fsz + "::" + bold + "::" + italic)
            elif (target == "main"):
               self.run_main.append(fnm + "::" + fsz + "::" + bold + "::" + italic)
            else:
               self.run_fn.append(fnm + "::" + fsz + "::" + bold + "::" + italic)
         #  Token is a word, so build up the text string
         else:
            if (text == ""):
               text = toks[0]
               runtext = toks[0]
            else:
               text = text + " " + toks[0]
               runtext = runtext + " " + toks[0]

         toks.pop(0)

      if (text != ""):
         if (target == "margin"):
            self.bullets_marg.append(text)
            self.run_marg.append(runtext)
         elif (target == "main"):
            self.bullets_main.append(text)
            self.run_main.append(runtext)
         else:
            self.footnotes.append(text)
            self.run_fn.append(runtext)

      return 0
